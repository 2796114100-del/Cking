# tools/file_tools.py
# 文件操作工具：支持 .txt .docx .pptx .xlsx 格式

import os


def read_file(file_path: str) -> str:
    """
    读取指定路径的文件内容。
    支持格式：.txt .docx .pptx .xlsx
    """
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.docx':
            return _read_docx(file_path)
        elif ext == '.pptx':
            return _read_pptx(file_path)
        elif ext == '.xlsx':
            return _read_xlsx(file_path)
        else:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
    except FileNotFoundError:
        return f"文件不存在：{file_path}"
    except Exception as e:
        return f"读取文件时出错了：{e}"


def write_file(file_path: str, content: str) -> str:
    """
    把内容写入指定文件。
    支持格式：.txt .docx .pptx .xlsx
    """
    try:
        folder = os.path.dirname(file_path)
        if folder and not os.path.exists(folder):
            os.makedirs(folder)

        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.docx':
            return _write_docx(file_path, content)
        elif ext == '.pptx':
            return _write_pptx(file_path, content)
        elif ext == '.xlsx':
            return _write_xlsx(file_path, content)
        else:
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            return f"写入成功！文件位置：{file_path}"
    except Exception as e:
        return f"写入文件时出错了：{e}"


# ==================== Word (.docx) ====================

def _read_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        parts.append("[表格]")
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts) if parts else "(空文档)"


def _write_docx(file_path: str, content: str) -> str:
    from docx import Document
    doc = Document()
    for line in content.split("\n"):
        doc.add_paragraph(line)
    doc.save(file_path)
    return f"Word文档已保存：{file_path}"


# ==================== PowerPoint (.pptx) ====================

def _read_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- 第{i}页 ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                parts.append("[表格]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    parts.append(row_text)
    return "\n".join(parts) if parts else "(空PPT)"


def _write_pptx(file_path: str, content: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    # 按 --- 分页
    slides = content.split("---")
    for slide_text in slides:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容布局
        lines = slide_text.split("\n")
        # 第一行作为标题
        if slide.shapes.title:
            slide.shapes.title.text = lines[0]
        # 其余作为内容
        if len(lines) > 1 and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            body.text = "\n".join(lines[1:])
    prs.save(file_path)
    return f"PPT已保存：{file_path}"


# ==================== Excel (.xlsx) ====================

def _read_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== 工作表：{sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts) if parts else "(空Excel)"


def _write_xlsx(file_path: str, content: str) -> str:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    rows = content.strip().split("\n")
    for r, row_text in enumerate(rows, 1):
        # 支持用 | 或 \t 分隔列
        if "\t" in row_text:
            cells = row_text.split("\t")
        elif "|" in row_text:
            cells = [c.strip() for c in row_text.split("|")]
        else:
            cells = [row_text]
        for c, val in enumerate(cells, 1):
            val = val.strip()
            # 尝试转数字
            try:
                val = float(val)
                if val == int(val):
                    val = int(val)
            except:
                pass
            ws.cell(row=r, column=c, value=val)
    wb.save(file_path)
    return f"Excel已保存：{file_path}"
